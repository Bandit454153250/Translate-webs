from flask import Flask, render_template, request, flash,redirect
from pptx import Presentation
from io import BytesIO
import os
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.secret_key = 'your-secret-key'  # 设置一个安全的密钥用于flash消息
app.config['UPLOAD_FOLDER'] = 'uploads'  # 设置上传文件夹
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 允许上传最大16MB的文件


@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST' and 'file' in request.files:
        file = request.files['file']
        if file.filename == '':
            flash('No selected file', 'error')
            return redirect(request.url)
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            
            # 使用pptx库读取PPT文件并提取文字
            prs = Presentation(filepath)
            text_items = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text_items.append(shape.text)
                    elif hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_items.append(run.text)

            # 假设我们有一个名为`display_text`的变量来保存提取的文本
            display_text = '\n'.join(text_items)

            # 渲染模板并传入提取的文本
            return render_template('web.html', display_text=display_text)

    return render_template('web.html')


def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'pptx'}
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


if __name__ == '__main__':
    app.run(debug=True)